# a tool for pasting figures into ppt.
# it extracts figures in each subfolder, rescales to the same width and height, and pastes them into one slide, i.e.
# creating one slide for each subfolder.
# supports: 1) only one level of subfolder; 2) a maximum of 8 figures in each subfolder.

# author: ligang zhang
# contact: ligang.zhang@btimaging.com
# version: 1
# creation time: 2023-02-02

# compile to .pyc files: python -m compileall main.py

import os
import glob
import argparse
import tifffile
import numpy as np
import PIL
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# add supported image types to the list
SUPPORT_EXT = ["png", "jpg", "jpeg", "gif"]
# converted type for tif or tiff images
CONVERT_EXT = SUPPORT_EXT[0]
SKIP_TIF_CONVERSION = False

def scale_image(im):
    """
    Scale an image to the range [0, 1]
    """
    dynamic_range = im.max() - im.min()

    if dynamic_range > 0:
        return (im - im.min()) / float(dynamic_range)
    else:
        return np.zeros_like(im)

def open_image(fn, rgb2gray=True, cast_long=True):
    im = None
    if os.path.splitext(fn)[1].lower() in ['.tif', '.tiff']:
        im = tifffile.imread(fn)
        if cast_long:
            im = im.astype(np.uint16)
    else:
        im = PIL.Image.open(fn)
        im = np.array(im)
    if im.ndim == 3 and rgb2gray:
        print('WARNING: 3-channel image. Using first channel.')
        im = np.ascontiguousarray(im[0, :, :])

    return im

def save_image(fnOut, im, scale=True):
    if im.dtype.type is np.uint8:
        if im.ndim == 2 and scale:
            im = (scale_image(im.astype(np.float64)) * 255).astype(np.uint8)
        pil_im = PIL.Image.fromarray(im)
        pil_im.save(fnOut)
    elif im.dtype.type is np.uint16:
        tifffile.imsave(fnOut, im)
    elif im.dtype.type is np.float32 and im.ndim == 2:
        im = (scale_image(im) * 255).astype(np.uint8)
        pil_im = PIL.Image.fromarray(im)
        pil_im.save(fnOut)
    else:
        print("ERROR: Unsupported file type. Type: %s  Dim: %d  Filename: %s" % (str(im.dtype.type), im.ndim, fnOut))
        sys.exit()

def main():

    # add arguments
    parser = argparse.ArgumentParser(description='A tool for pasting figures to PPT')
    current_dir = os.path.dirname(os.path.abspath(__file__))
    parser.add_argument('-i', type=str, default=os.path.join(current_dir, "input"), help = "input directory, default: %(default)s")
    parser.add_argument('-o', type=str, default=os.path.join(current_dir, "output"), help = "output directory, default: %(default)s")
    parser.add_argument('-n', type=str, default="ppt_name", help = "ppt name, default: %(default)s")
    parser.add_argument('-r', type=int, default=1, help = "keep image ratio: default: %(default)s")
    args = parser.parse_args()
    input_path = args.i
    output_path = args.o

    # saved ppt name
    ppt_name = args.n

    # keep original image ratio, or use same width and height.
    keep_im_ratio = args.r > 0

    # indentation for second level of output message
    spacing = ' '
    num_spacing = 4

    if not os.path.exists(input_path):
        print (f"ERROR: {input_path} does not exist!")
        sys.exit(-1)

    if not os.path.exists(output_path):
        os.makedirs(output_path)
    else:
        print (f"Warning: {output_path} already exists.")
        print ("         existing ppt may be overwritten!")

    if not SKIP_TIF_CONVERSION:
        # convert tif to png if tifs exist
        print(f"Converting tif and tiff to {CONVERT_EXT}...") 
        for root, _, files in os.walk(input_path):
            for file in files:
                if file.lower().endswith(('.tiff', '.tif')):
                    fn = os.path.join(root, file)
                    
                    im_name = os.path.splitext(os.path.split(fn)[1])[0]
                    im_out = os.path.join(root, im_name + f".{CONVERT_EXT}")
                    if os.path.exists(im_out):
                        print (f"{spacing*num_spacing}Warning: {im_out} already exists. Conversion skipped!")
                        continue
                    im = open_image(fn).astype(np.float32)
                    save_image(im_out, im)
        print("Convertion done!")

    prs = Presentation()
    # default slide width
    #prs.slide_height = 9144000
    # slide height @ 4:3
    #prs.slide_height = 6858000
    # slide height @ 16:9
    prs.slide_height = 5143500
    SH = prs.slide_height
    SW = prs.slide_width

    # set title position for the 2nd and following slides
    title_top   = int(SH * 0.03)
    title_height = int(SH * 0.15)
    title_width = int(SW * 0.85)

    # set figure size and position
    slide_top_14 =  int(SH * 0.25)
    slide_top_58 =  int(SH * 0.18)
    fig_height_12   = int(SH * 0.65)
    fig_height_3   = int(SH * 0.5)
    fig_height_48   = int(SH * 0.37)
    figure_spacing   = int(SH * 0.04)

    # set textbox position for slide number
    number_left = int(SW * 0.95)
    number_top = int(SH * 0.95)
    number_width = int(SW * 0.05)
    number_height = int(SW * 0.03)

    max_figures_per_row = 4

    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # set title
    title = slide.shapes.title
    title.text = "add title here"

     # for images in all subdirectories
    print(f"Pasting figures to PPT...")
    slide_number = 0
    for folder_name in os.listdir(input_path):
        folder_path = os.path.join(input_path, folder_name)
        print (f"{spacing*num_spacing}folder: {folder_path}")
        if os.path.isdir(folder_path):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide_number += 1
            
            # set title
            title = slide.shapes.title
            title.text = folder_name
            title.top = title_top
            title.height = title_height
            title.width = title_width
            # force title box at center
            title.left = (SW - title.width)//2

            # control text size and alignment
            title_para = title.text_frame.paragraphs[0]
            title_para.font.size = Pt(40)
            title_para.alignment = PP_ALIGN.CENTER

            # get path of images with supported extensions 
            fns = []
            for ext in SUPPORT_EXT:
                fns.extend(glob.glob(os.path.join(folder_path, f"*.{ext}")))

            num_figures = len(fns)
            if num_figures == 0:
                print (f"{spacing*num_spacing}Warning: no {SUPPORT_EXT} images found in {folder_path}!")
                continue
            
            for e, fn in enumerate(fns):
                im = open_image(fn, rgb2gray=False).astype(np.float32)
                im_ratio = im.shape[1]/im.shape[0]

                # set layouts for different numbers of figures
                if num_figures <= 4:
                    fig_top = slide_top_14
                    if num_figures <= 2:
                        fig_height = fig_height_12
                    elif num_figures == 3:
                        fig_height = fig_height_3
                    elif num_figures == 4:
                        fig_height = fig_height_48

                    # calcualate the left location for the first figure
                    if e==0:
                        fig_left = (SW - fig_height * num_figures - figure_spacing * (num_figures-1))//2
                    else:
                        fig_left += fig_height + figure_spacing
                    
                    if keep_im_ratio:
                        # keep original image ratio
                        if im_ratio > 1.0:
                            fig_width = fig_height
                            fig_height = int(fig_width / im_ratio)
                        else:
                            fig_width = int(fig_height * im_ratio)
                    else:
                        # force the same width and height
                        fig_width = fig_height
                elif num_figures <= 8:
                    # spread to two rows
                    fig_top = slide_top_58
                    fig_height = fig_height_48

                    if e == 0 or e == 4:
                        fig_left = (SW - fig_height * max_figures_per_row - figure_spacing * (max_figures_per_row-1))//2
                    else:
                        fig_left += fig_height + figure_spacing

                    if e >= 4:
                        fig_top = fig_top + fig_height + figure_spacing

                    if keep_im_ratio:
                        # keep original image ratio
                        if im_ratio > 1.0:
                            fig_width = fig_height
                            fig_height = int(fig_width / im_ratio)
                        else:
                            fig_width = int(fig_height * im_ratio)
                    else:
                        # force the same width and height
                        fig_width = fig_height
                else:
                    print(f"{spacing*num_spacing}ERROR: support only a maximum of 8 images per slide!")
                    sys.exit()

                # add figure to ppt
                slide.shapes.add_picture(fn, fig_left, fig_top, fig_width, fig_height)

                # set textbox position
                text_left = fig_left
                text_top = fig_top + fig_height - int(SH * 0.01)
                text_width = fig_height
                text_height = figure_spacing

                # add image name
                txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                im_name = os.path.split(fn)[1].strip()
                p.text = im_name
                p.font.size = Pt(13)

            # add slide number
            txBox = slide.shapes.add_textbox(number_left, number_top, number_width, number_height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_number)
            p.font.size = Pt(13)
    
    # save ppt
    ppt_path = os.path.join(output_path, f"{ppt_name}.pptx")
    prs.save(ppt_path)
    print(f"Done! PPT saved at {ppt_path}.")

if __name__ == '__main__':
    main()